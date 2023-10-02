"""_summary_
"""
from PyQt5.QtWidgets import QFileDialog, QApplication
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from PIL import Image
from io import BytesIO
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPDF
import cairosvg
import os
import sys
import tempfile


def printProgressBar(iteration, total, prefix='', suffix='', decimals=1,
                     length=100, fill='█', printEnd="\r"):
    """
    Call in a loop to create terminal progress bar.

    Parameters
    ----------
    iteration : `int`
        Current iteration.
    total : `int`
        Total iterations.
    prefix : `str`, optional
        Prefix string.
    suffix : `str`, optional
        Suffix string.
    decimals : `int`, optional
        Positive number of decimals in percent.
    length : `int`, optional
        Character length of the bar.
    fill : `str`, optional
        Bar fill character.
    printEnd : `str`, optional
        End character (e.g., "\\r", "\\r\\n").
    """
    percent = ("{0:." + str(decimals) + "f}").format(100 * (iteration
                                                            / float(total)))
    filledLength = int(length * iteration // total)
    bar = fill * filledLength + '-' * (length - filledLength)
    print(f'\r{prefix} |{bar}| {percent}% {suffix}', end=printEnd)
    if iteration == total:
        print()


def insert_figures(folder, x, y, gap_x, img_width):
    """
    Extract svg files from folder and insert into a row in ppt slide.

    Bit janky. Starts at x, y. Converts svg to png insert into row. Each item
    of row moves over by img_width + gap_x
    Parameters
    ----------
    folder : str
        filepath to the folder containing images
    x : float
        float representing inches value of position to start adding photos
    y : float
        float representing inches value of position to start adding photos
    gap_x : float
        float representing inches value of gap between photos
    img_width : float
        float representing inches value of photo width
    """
    # Iterate over the image files in the Results folder
    for idx, image_file in enumerate(os.listdir(folder)):
        if image_file.endswith(".svg"):
            img_path = os.path.join(folder, image_file)
            left_offset = x + idx % 3 * (img_width + gap_x)

            # Render SVG to a temporary PNG file using cairosvg
            temp_png = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            temp_png.close()  # Close the file so that it can be opened for writing

            cairosvg.svg2png(url=img_path, write_to=temp_png.name)

            # Add the image to the slide
            slide.shapes.add_picture(temp_png.name,
                                     Inches(left_offset), Inches(y),
                                     width=Inches(img_width))
            # Remove the temporary PNG file
            os.unlink(temp_png.name)


app = QApplication(sys.argv)  # noqa
# Path to the main folder
main_folder = QFileDialog.getExistingDirectory(None, 'Select Directory of study')

print('Generating Summary Slides...')

# Initialize a new PowerPoint presentation
prs = Presentation()
# Set the slide size to 16:9 (widescreen) aspect ratio
prs.slide_width = Inches(13.333)  # 16:9 width in inches
prs.slide_height = Inches(7.5)    # 16:9 height in inches

# Walk through the main folder and its subdirectories
#folders = sorted([folder for folder in os.listdir(main_folder) if os.path.isdir(os.path.join(main_folder, folder))])

folderpaths = [folder for folder in os.listdir(main_folder) if os.path.isdir(os.path.join(main_folder, folder))]
folderpaths = [os.path.join(main_folder, f) for f in folderpaths] # add path to each file
folderpaths.sort(key=lambda x: os.path.getmtime(x))

# Process folders in pairs
for i in range(0, len(folderpaths), 2):
    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0)
    img_width = 4.35
    img_height = Inches(3.25)

    # Identify first folder and expt
    folder1 = folderpaths[i]
    expt_name1 = os.path.basename(folder1)
    results_fol1 = os.path.join(folder1, "Results")

    # Add the figures to top half of slide
    insert_figures(results_fol1, 0, 0.6, 0.15, img_width)

    # Create textbox for expt name
    expt1_label = slide.shapes.add_textbox(Inches(0.25), Inches(0.15),
                                           Inches(9), Inches(0.5))
    expt1_label.text = expt_name1
    expt1_label.text_frame.paragraphs[0].font.size = Pt(18)
    expt1_label.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Check if there is a paired folder available
    if i + 1 < len(folderpaths):
        # Get folder path and expt name
        folder2 = folderpaths[i+1]
        expt_name2 = os.path.basename(folder2)
        results_fol2 = os.path.join(folder2, "Results")

        # Add the figures to bottom half of slide
        insert_figures(results_fol2, 0, 4.25, 0.15, img_width)

        # Create textbox for expt name
        expt2_label = slide.shapes.add_textbox(Inches(0.25), Inches(4),
                                               Inches(9), Inches(0.5))
        expt2_label.text = expt_name2
        expt2_label.text_frame.paragraphs[0].font.size = Pt(18)
        expt2_label.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    printProgressBar(i, len(folderpaths))
printProgressBar(1, 1)
# Save the PowerPoint presentation
prs.save(os.path.join(main_folder, "summary_slides.pptx"))
print('Finished!')
